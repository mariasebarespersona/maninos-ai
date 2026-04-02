"""
Generate a single-slide professional PPT showcasing Maninos AI automations.
Navy (#1E3A5F) + Gold (#C9A227) theme — elegant, illustrative layout.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Palette ---
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
NAVY_DARK = RGBColor(0x0F, 0x1F, 0x33)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GOLD_LIGHT = RGBColor(0xE8, 0xC8, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MEDIUM_GRAY = RGBColor(0x8A, 0x94, 0xA6)
SOFT_WHITE = RGBColor(0xFA, 0xFA, 0xFC)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Helper functions ──────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def set_text(shape, text, font_size=12, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_text(tf, text, font_size=10, color=WHITE, indent=0, bold=False, font_name="Calibri", space_before=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.level = indent
    p.space_before = space_before
    p.alignment = PP_ALIGN.LEFT
    return p

# ── 1. BACKGROUND ────────────────────────────────────────────────
add_rect(slide, Inches(0), Inches(0), Inches(16), Inches(9), fill=NAVY_DARK)

# Subtle accent stripe at top
add_rect(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), fill=GOLD)

# ── 2. HEADER ────────────────────────────────────────────────────
# Title area
title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(10), Inches(0.65))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT

run1 = p.add_run()
run1.text = "MANINOS AI"
run1.font.size = Pt(32)
run1.font.color.rgb = GOLD
run1.font.bold = True
run1.font.name = "Calibri"

run2 = p.add_run()
run2.text = "  |  Plataforma Inteligente de Automatizaciones"
run2.font.size = Pt(18)
run2.font.color.rgb = MEDIUM_GRAY
run2.font.bold = False
run2.font.name = "Calibri"

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(14), Inches(0.4))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Ecosistema completo de IA para adquisicion, renovacion, venta, financiacion y contabilidad de propiedades"
p.font.size = Pt(11)
p.font.color.rgb = MEDIUM_GRAY
p.font.name = "Calibri"
p.alignment = PP_ALIGN.LEFT

# Gold divider line
add_rect(slide, Inches(0.8), Inches(1.45), Inches(14.4), Inches(0.02), fill=GOLD)

# ── Stats bar (right side of header) ─────────────────────────────
stats = [("5", "Agentes IA"), ("7+", "Fuentes Datos"), ("58+", "Migraciones"), ("3", "Portales")]
for i, (num, label) in enumerate(stats):
    x = Inches(11.5 + i * 1.15)
    nb = slide.shapes.add_textbox(x, Inches(0.35), Inches(1.0), Inches(0.35))
    tf = nb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.size = Pt(22)
    r.font.color.rgb = GOLD
    r.font.bold = True
    r.font.name = "Calibri"

    lb = slide.shapes.add_textbox(x, Inches(0.68), Inches(1.0), Inches(0.3))
    tf2 = lb.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.text = label
    p2.font.size = Pt(8)
    p2.font.color.rgb = MEDIUM_GRAY
    p2.font.name = "Calibri"

# ── 3. AUTOMATION CARDS ──────────────────────────────────────────

cards = [
    {
        "icon": "🔍",
        "title": "CASAS DEL MERCADO",
        "subtitle": "Scraping Inteligente Multi-Fuente",
        "bullets": [
            "7 fuentes: MHVillage, MobileHome.net, Facebook, VMF Homes, 21st Mortgage, Zillow, Craigslist",
            "Filtros automaticos: rango $5K-$80K, radio 200mi Houston/Dallas, regla del 60%",
            "Playwright + GPT-4 Vision: extraccion de screenshots de Facebook Marketplace",
            "Deduplicacion inteligente por normalizacion de direcciones entre fuentes",
            "Auto-refresco cada 6h (APIs partner) + deteccion de down payment vs precio total",
            "150+ listings/fuente desde JSON APIs directas sin browser",
        ],
        "accent": RGBColor(0x4E, 0xA8, 0xDE),  # blue
    },
    {
        "icon": "📄",
        "title": "DOCUMENTOS AUTOMATICOS",
        "subtitle": "Auto-Generacion de Documentos de Compra",
        "bullets": [
            "Bill of Sale, Title Transfer (TDHCA SOL), Depositos, Contratos RTO",
            "Auto-relleno con datos del cliente, propiedad, HUD, VIN, precio",
            "PDFs branded (navy/gold) via ReportLab + almacenamiento en Supabase",
            "Generacion automatica al confirmar pago (Contado y RTO)",
            "Monitoreo TDHCA: alerta automatica cuando cambia nombre de titulo",
            "Pagares, estados de cuenta y reportes financieros en PDF",
        ],
        "accent": RGBColor(0x6C, 0xBF, 0x84),  # green
    },
    {
        "icon": "🏗️",
        "title": "RENOVACIONES + VOZ",
        "subtitle": "Gestion IA con Comandos de Voz",
        "bullets": [
            "Checklist de 19 items + items personalizados por propiedad",
            "Agente de Voz (Whisper + GPT-4): dictado manos libres en campo",
            "Fotos → IA: GPT-4 Vision analiza fotos y auto-rellena checklist",
            "Agente Costos: estimacion inteligente materiales + mano de obra",
            "Importacion de reportes de evaluacion (14 items inspeccion)",
            "Flujo de aprobacion gerencial antes de iniciar obra",
        ],
        "accent": RGBColor(0xE8, 0x9B, 0x3E),  # orange
    },
    {
        "icon": "📊",
        "title": "CONTABILIDAD IA",
        "subtitle": "Software Financiero con Agente Inteligente",
        "bullets": [
            "Conciliacion bancaria automatica + multi-cuenta (Zelle, Stripe, transferencia)",
            "IA categoriza transacciones y reconoce patrones recurrentes",
            "Estados financieros: Balance, P&L, Flujo de Caja + presupuesto vs real",
            "P&L por propiedad y por yard (Houston/Conroe/Dallas)",
            "Journal completo con audit trail + gastos recurrentes automaticos",
            "Exportacion CSV + dashboard con KPIs financieros en tiempo real",
        ],
        "accent": RGBColor(0xAF, 0x7A, 0xC5),  # purple
    },
    {
        "icon": "💰",
        "title": "PORTAL CAPITAL",
        "subtitle": "Inversores, RTO y Financiacion Inteligente",
        "bullets": [
            "6 fases: Analisis → Adquirir → Firmar → Gestionar → Reportes → Fondear",
            "Analisis financiero: ROI, breakeven, riesgo, precios sugeridos RTO",
            "KYC verificacion (documentos + selfie) + scoring de aplicaciones",
            "Tracking pagos mensuales + alertas automaticas (3d antes, dia de, 1d despues)",
            "Gestion de inversores: capital flows, ciclos, notas promisorias",
            "Comisiones auto: $1,500 contado / $1,000 RTO con split 50/50 finder/closer",
        ],
        "accent": RGBColor(0xDE, 0x6B, 0x6B),  # red
    },
]

# Layout: 3 cards top row, 2 cards bottom row (centered)
card_w = Inches(4.65)
card_h = Inches(3.15)
margin_x = Inches(0.8)
gap = Inches(0.12)
y_row1 = Inches(1.65)
y_row2 = Inches(5.0)

positions = []
# Row 1: 3 cards
for i in range(3):
    x = margin_x + i * (card_w + gap)
    positions.append((x, y_row1))
# Row 2: 2 cards centered
row2_total = 2 * card_w + gap
row2_start = (Inches(16) - row2_total) / 2
for i in range(2):
    x = row2_start + i * (card_w + gap)
    positions.append((x, y_row2))

for idx, card in enumerate(cards):
    cx, cy = positions[idx]
    accent = card["accent"]

    # Card background
    bg = add_rounded_rect(slide, cx, cy, card_w, card_h, fill=NAVY, border_color=RGBColor(0x2A, 0x4A, 0x6F), border_width=Pt(1))

    # Accent bar at top of card
    add_rect(slide, cx + Inches(0.0), cy, card_w, Inches(0.04), fill=accent)

    # Icon + Title
    header_box = slide.shapes.add_textbox(cx + Inches(0.25), cy + Inches(0.15), card_w - Inches(0.5), Inches(0.4))
    tf = header_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    r_icon = p.add_run()
    r_icon.text = card["icon"] + "  "
    r_icon.font.size = Pt(16)

    r_title = p.add_run()
    r_title.text = card["title"]
    r_title.font.size = Pt(14)
    r_title.font.color.rgb = WHITE
    r_title.font.bold = True
    r_title.font.name = "Calibri"

    # Subtitle
    sub = slide.shapes.add_textbox(cx + Inches(0.25), cy + Inches(0.52), card_w - Inches(0.5), Inches(0.25))
    tf = sub.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = card["subtitle"]
    p.font.size = Pt(9)
    p.font.color.rgb = accent
    p.font.bold = True
    p.font.name = "Calibri"

    # Thin divider
    add_rect(slide, cx + Inches(0.25), cy + Inches(0.78), card_w - Inches(0.5), Inches(0.01), fill=RGBColor(0x2A, 0x4A, 0x6F))

    # Bullets
    bullet_box = slide.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.82), card_w - Inches(0.4), card_h - Inches(1.0))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for bi, bullet in enumerate(card["bullets"]):
        if bi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▸  " + bullet
        p.font.size = Pt(8)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Calibri"
        p.space_before = Pt(2)
        p.space_after = Pt(1)

# ── 4. TECH STACK FOOTER ────────────────────────────────────────
add_rect(slide, Inches(0), Inches(8.4), Inches(16), Inches(0.6), fill=NAVY_DARK)
add_rect(slide, Inches(0), Inches(8.38), Inches(16), Inches(0.015), fill=GOLD)

# Tech stack text
tech_box = slide.shapes.add_textbox(Inches(0.8), Inches(8.45), Inches(12), Inches(0.45))
tf = tech_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT

techs = [
    ("Stack:  ", MEDIUM_GRAY, False, Pt(8)),
    ("FastAPI", GOLD_LIGHT, True, Pt(8)),
    ("  ·  ", MEDIUM_GRAY, False, Pt(8)),
    ("Next.js 14", GOLD_LIGHT, True, Pt(8)),
    ("  ·  ", MEDIUM_GRAY, False, Pt(8)),
    ("GPT-4o + Whisper", GOLD_LIGHT, True, Pt(8)),
    ("  ·  ", MEDIUM_GRAY, False, Pt(8)),
    ("Supabase", GOLD_LIGHT, True, Pt(8)),
    ("  ·  ", MEDIUM_GRAY, False, Pt(8)),
    ("Stripe", GOLD_LIGHT, True, Pt(8)),
    ("  ·  ", MEDIUM_GRAY, False, Pt(8)),
    ("Playwright", GOLD_LIGHT, True, Pt(8)),
    ("  ·  ", MEDIUM_GRAY, False, Pt(8)),
    ("Resend", GOLD_LIGHT, True, Pt(8)),
    ("  ·  ", MEDIUM_GRAY, False, Pt(8)),
    ("Railway + Vercel", GOLD_LIGHT, True, Pt(8)),
]
for text, color, bold, size in techs:
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Calibri"

# Notifications/scheduling extra info
extra_box = slide.shapes.add_textbox(Inches(0.8), Inches(8.7), Inches(14), Inches(0.25))
tf = extra_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
extras = "Emails automaticos (bienvenida, pago, recordatorios RTO, alertas morosidad, reviews, referidos)  ·  5 jobs programados (APScheduler)  ·  3 portales: Homes, Capital, Clientes"
p.text = extras
p.font.size = Pt(7)
p.font.color.rgb = MEDIUM_GRAY
p.font.name = "Calibri"

# Brand mark right side
brand = slide.shapes.add_textbox(Inches(13.5), Inches(8.48), Inches(2.2), Inches(0.35))
tf = brand.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
r = p.add_run()
r.text = "RAMA AI"
r.font.size = Pt(12)
r.font.color.rgb = GOLD
r.font.bold = True
r.font.name = "Calibri"

# Save
output_path = "/Users/mariasebares/Documents/RAMA_AI/maninos-ai/Maninos_AI_Automations.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
