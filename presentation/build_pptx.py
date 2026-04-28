"""
Build Maninos OS pitch deck as PPTX, mirroring the HTML deck.
Editorial tone, navy/gold palette, only verified numbers from the repo.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os

# Palette
BG = RGBColor(0x0D, 0x16, 0x23)
BG_2 = RGBColor(0x13, 0x1F, 0x30)
SURFACE = RGBColor(0x1A, 0x26, 0x38)
LINE = RGBColor(0x2A, 0x36, 0x50)
LINE_SOFT = RGBColor(0x1F, 0x2A, 0x3E)
INK = RGBColor(0xEE, 0xF0, 0xF4)
INK2 = RGBColor(0xC7, 0xCD, 0xD8)
INK3 = RGBColor(0x8C, 0x95, 0xA8)
INK4 = RGBColor(0x5A, 0x64, 0x78)
GOLD = RGBColor(0xC9, 0xA0, 0x47)
GOLD_SOFT = RGBColor(0xE6, 0xCF, 0x95)

SERIF = "Georgia"          # closest universally available serif (Fraunces fallback)
SANS = "Calibri"           # universal
MONO = "Consolas"          # universal mono

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Outer margin
MARGIN_X = Inches(0.6)
MARGIN_Y = Inches(0.45)
CONTENT_W = SLIDE_W - 2 * MARGIN_X


def _set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *, font=SANS, size=14, color=INK,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.3, letter_spacing=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        runs = [(text, {})]
    else:
        runs = text  # list of (str, dict)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for i, (chunk, opts) in enumerate(runs):
        r = p.add_run() if i > 0 else p.runs[0] if p.runs else p.add_run()
        # workaround: ensure exactly one run
        r.text = chunk
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.italic = opts.get("italic", italic)
        r.font.color.rgb = opts.get("color", color)
        # letter spacing
        ls = opts.get("letter_spacing", letter_spacing)
        if ls:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(ls * 100)))
    return tb


def add_paragraph(tf, text, *, font=SANS, size=14, color=INK,
                  bold=False, italic=False, align=PP_ALIGN.LEFT,
                  line_spacing=1.3, letter_spacing=0, space_after=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    if space_after:
        p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if letter_spacing:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing * 100)))
    return p


def add_eyebrow(slide, x, y, text, w=Inches(8)):
    return add_text(slide, x, y, w, Inches(0.3), text,
                    size=9, color=GOLD, bold=True, letter_spacing=2.5)


def add_rule(slide, x, y, w, color=LINE, thickness=Pt(0.75)):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = thickness
    return line


def add_vrule(slide, x, y, h, color=LINE_SOFT, thickness=Pt(0.75)):
    line = slide.shapes.add_connector(1, x, y, x, y + h)
    line.line.color.rgb = color
    line.line.width = thickness
    return line


def add_footer(slide, left_text, right_text):
    add_text(slide, MARGIN_X, SLIDE_H - Inches(0.4), CONTENT_W, Inches(0.25),
             left_text, size=8, color=INK4, letter_spacing=2, bold=True)
    add_text(slide, MARGIN_X, SLIDE_H - Inches(0.4), CONTENT_W, Inches(0.25),
             right_text, size=8, color=INK4, letter_spacing=2, bold=True,
             align=PP_ALIGN.RIGHT)


def slide_blank(prs):
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)
    _set_bg(s)
    return s


# ============== SLIDES ==============

def slide_title(prs):
    s = slide_blank(prs)

    add_text(s, MARGIN_X, Inches(0.6), CONTENT_W, Inches(0.3),
             "MANINOS OS    ·    DOCUMENTO DESCRIPTIVO    ·    VERSIÓN 2026.04",
             size=9, color=INK3, letter_spacing=3, bold=True)

    add_text(s, MARGIN_X, Inches(2.0), CONTENT_W, Inches(2.5),
             "Una plataforma operativa\npara vender mobile homes,",
             font=SERIF, size=54, color=INK, line_spacing=1.05)

    add_text(s, MARGIN_X, Inches(3.7), CONTENT_W, Inches(1.0),
             "de extremo a extremo.",
             font=SERIF, size=54, color=GOLD, italic=True, line_spacing=1.05)

    add_text(s, MARGIN_X, Inches(5.0), Inches(8), Inches(1),
             "Cómo está construida, qué resuelve, y por qué su arquitectura admite\notros activos físicos con el mismo patrón de negocio.",
             font=SERIF, size=18, color=INK2, line_spacing=1.4)

    # bottom meta
    add_text(s, MARGIN_X, SLIDE_H - Inches(0.7), Inches(8), Inches(0.3),
             "MANINOS HOMES LLC  ·  MANINOS CAPITAL LLC",
             size=8, color=INK4, letter_spacing=2.5, bold=True)
    add_text(s, MARGIN_X, SLIDE_H - Inches(0.7), CONTENT_W, Inches(0.3),
             "TEXAS — 2026",
             size=8, color=INK4, letter_spacing=2.5, bold=True, align=PP_ALIGN.RIGHT)


def slide_definicion(prs):
    s = slide_blank(prs)
    add_eyebrow(s, MARGIN_X, MARGIN_Y, "DEFINICIÓN")
    add_text(s, MARGIN_X, MARGIN_Y + Inches(0.3), CONTENT_W, Inches(0.9),
             "Qué es Maninos OS",
             font=SERIF, size=36, color=INK)

    lead = (
        "Es la plataforma de software propietaria que opera Maninos Homes LLC y "
        "Maninos Capital LLC. Cubre el ciclo completo de un activo: localización en "
        "el mercado, compra, renovación, listado, venta, financiamiento propio "
        "(RTO o contado), gestión de pagos recurrentes, transferencia de título y "
        "reporte a inversores."
    )
    add_text(s, MARGIN_X, Inches(2.0), Inches(11.5), Inches(2),
             lead, font=SERIF, size=17, color=INK, line_spacing=1.4)

    add_rule(s, MARGIN_X, Inches(4.1), CONTENT_W)

    # two col
    col_w = (CONTENT_W - Inches(0.6)) / 2
    col1_x = MARGIN_X
    col2_x = MARGIN_X + col_w + Inches(0.6)

    add_text(s, col1_x, Inches(4.4), col_w, Inches(0.3),
             "LO QUE REEMPLAZA", size=8, color=GOLD, bold=True, letter_spacing=2)

    items_left = [
        "Hojas de cálculo y grupos de WhatsApp para inventario y renovación",
        "Un CRM genérico para clientes y leads",
        "Un sistema de contabilidad externo",
        "Un proceso manual de transferencia de título estatal",
        "Una hoja paralela para tracking de inversores y mora",
    ]
    add_text(s, col1_x, Inches(4.75), col_w, Inches(0.3),
             "—  " + items_left[0], size=11, color=INK2, line_spacing=1.5)
    y = Inches(5.1)
    for itm in items_left[1:]:
        add_text(s, col1_x, y, col_w, Inches(0.3),
                 "—  " + itm, size=11, color=INK2, line_spacing=1.5)
        y += Inches(0.32)

    add_text(s, col2_x, Inches(4.4), col_w, Inches(0.3),
             "LO QUE AÑADE QUE NO EXISTE EN OTRA PARTE",
             size=8, color=GOLD, bold=True, letter_spacing=2)
    items_right = [
        "Seis agentes de IA especializados, integrados en los flujos de trabajo",
        "Un módulo de financiamiento propio (RTO) con su propia contabilidad",
        "Un portal público para el cliente final con simulador, KYC, firma y autoservicio de pagos",
        "Un módulo separado para gestión de capital de inversores privados",
    ]
    y = Inches(4.75)
    for itm in items_right:
        add_text(s, col2_x, y, col_w, Inches(0.5),
                 "—  " + itm, size=11, color=INK2, line_spacing=1.5)
        y += Inches(0.4)

    add_footer(s, "01  ·  DEFINICIÓN", "MANINOS OS")


def slide_arquitectura(prs):
    s = slide_blank(prs)

    # title row
    add_text(s, MARGIN_X, MARGIN_Y, Inches(8), Inches(0.6),
             "La forma del sistema", font=SERIF, size=32, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "ARQUITECTURA", size=8, color=INK3, bold=True, letter_spacing=2.5,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.45), Inches(11.5), Inches(0.9),
             "Tres portales que comparten una sola base de datos Postgres. No hay sincronización ni jobs de réplica entre ellos: lo que ocurre en operaciones se ve al instante en finanzas y en la cuenta del cliente.",
             font=SERIF, size=15, color=INK, line_spacing=1.4)

    # 3 columns
    cw = (CONTENT_W - Inches(0.8)) / 3
    base_y = Inches(2.9)

    portals = [
        ("PORTAL 01", "Homes", "OPERACIONES INTERNAS",
         "Para el equipo de Maninos. Inventario, renovación, fotos, contabilidad, comisiones, títulos, dashboards. Acceso por rol: treasury, operations, yard_manager, admin, comprador, vendedor."),
        ("PORTAL 02", "Capital", "FINANCIAMIENTO E INVERSORES",
         "Para underwriters y partners financieros. Solicitudes RTO, contratos de arrendamiento con opción a compra, pagos mensuales, mora, KYC, capital de inversores, promissory notes, reportes mensuales."),
        ("PORTAL 03", "Clientes", "CARA PÚBLICA Y AUTOSERVICIO",
         "Para el comprador. Catálogo público, simulador de RTO, solicitud de crédito, KYC con subida de documentos, firma electrónica del contrato, estado de cuenta y reporte de pagos."),
    ]
    for i, (eb, title, sub, body) in enumerate(portals):
        x = MARGIN_X + i * (cw + Inches(0.4))
        add_rule(s, x, base_y, cw)
        add_text(s, x, base_y + Inches(0.15), cw, Inches(0.3),
                 eb, size=8, color=GOLD, bold=True, letter_spacing=2)
        add_text(s, x, base_y + Inches(0.5), cw, Inches(0.6),
                 title, font=SERIF, size=22, color=INK)
        add_text(s, x, base_y + Inches(1.05), cw, Inches(0.3),
                 sub, size=8, color=INK3, bold=True, letter_spacing=2)
        add_text(s, x, base_y + Inches(1.4), cw, Inches(2),
                 body, size=10.5, color=INK2, line_spacing=1.5)

    # KPI row
    add_rule(s, MARGIN_X, Inches(5.95), CONTENT_W)
    add_rule(s, MARGIN_X, Inches(6.65), CONTENT_W)
    kpis = [
        ("18", "PÁGINAS EN HOMES"),
        ("19", "PÁGINAS EN CAPITAL"),
        ("16", "PÁGINAS EN CLIENTES"),
        ("58", "TABLAS EN POSTGRES"),
        ("91", "MIGRACIONES SQL"),
    ]
    kw = CONTENT_W / 5
    for i, (num, lbl) in enumerate(kpis):
        x = MARGIN_X + i * kw
        add_text(s, x + Inches(0.05), Inches(6.05), kw, Inches(0.45),
                 num, font=SERIF, size=24, color=INK)
        add_text(s, x + Inches(0.05), Inches(6.45), kw, Inches(0.2),
                 lbl, size=7, color=INK3, bold=True, letter_spacing=1.5)
        if i < 4:
            add_vrule(s, x + kw, Inches(6.0), Inches(0.65))

    add_text(s, MARGIN_X, Inches(6.85), CONTENT_W, Inches(0.3),
             "Cifras tomadas directamente del repositorio de Maninos a fecha de hoy.",
             size=9, color=INK3, italic=True)

    add_footer(s, "02  ·  ARQUITECTURA", "TRES PORTALES  ·  UNA BASE DE DATOS")


def slide_section_index(prs, num, title, desc):
    s = slide_blank(prs)
    _set_bg(s, BG_2)
    add_text(s, MARGIN_X, Inches(2.5), CONTENT_W, Inches(0.4),
             num, font=SERIF, size=12, color=GOLD, bold=True, letter_spacing=2)
    add_text(s, MARGIN_X, Inches(3.0), Inches(10), Inches(2),
             title, font=SERIF, size=48, color=INK, line_spacing=1.05)
    add_text(s, MARGIN_X, Inches(5.0), Inches(8), Inches(1.5),
             desc, font=SERIF, size=15, color=INK3, italic=True, line_spacing=1.5)


def slide_homes_portal(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(8), Inches(0.6),
             "Portal Homes", font=SERIF, size=32, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "OPERACIONES INTERNAS", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    # left sidebar label
    add_text(s, MARGIN_X, Inches(1.6), Inches(2.5), Inches(0.3),
             "/HOMES", font=MONO, size=8, color=INK4, letter_spacing=1.5, bold=True)
    add_text(s, MARGIN_X, Inches(1.95), Inches(2.5), Inches(1),
             "El día a día del staff",
             font=SERIF, size=15, color=GOLD, italic=True, line_spacing=1.3)

    # right column
    rx = MARGIN_X + Inches(2.9)
    rw = CONTENT_W - Inches(2.9)
    add_text(s, rx, Inches(1.6), rw, Inches(1),
             "El staff de Maninos pasa la jornada en este portal. Accede por /homes/* tras autenticarse en Supabase con su rol asignado. La barra lateral muestra solo las secciones a las que tiene acceso, con un badge de pendientes (transferencias de título, órdenes de pago, aprobaciones de renovación).",
             size=11, color=INK2, line_spacing=1.5)

    # table
    rows = [
        ("Dashboard", "KPIs, mapa de Texas, gráfico de actividad de los últimos 6 meses, pipeline por estado de propiedad."),
        ("Properties", "Listado con filtros por estado (purchased, published, reserved, renovating, sold), búsqueda por dirección/código. Subrutas para edición, fotos, checklist “Revisar Casa” y wizard de renovación."),
        ("Transfers", "Títulos en curso. Widget del scheduler diario contra TDHCA + modal para subida manual de títulos legacy. Cada serial es un link directo a la base TDHCA."),
        ("Sales", "Ventas activas, pagos parciales, comisiones (finder + closer), generación de Bill of Sale en PDF, redirección a Capital cuando es RTO."),
        ("Clients", "CRM interno solo para clientes contado."),
        ("Accounting", "11 pestañas: overview, transactions, invoices, statements, chart of accounts, properties, banks, budget, recurring, audit, estado de cuenta."),
        ("Commissions", "Comisiones por empleado. Treasury y Admin ven todas; otros roles solo las propias."),
        ("Market", "Listings externos scrapeados de fuentes públicas."),
    ]
    table_y = Inches(2.65)
    col1_w = Inches(1.8)
    col2_w = rw - col1_w - Inches(0.2)
    # header
    add_text(s, rx, table_y, col1_w, Inches(0.25),
             "SECCIÓN", size=7, color=INK3, bold=True, letter_spacing=1.8)
    add_text(s, rx + col1_w, table_y, col2_w, Inches(0.25),
             "FUNCIÓN", size=7, color=INK3, bold=True, letter_spacing=1.8)
    add_rule(s, rx, table_y + Inches(0.28), rw)
    y = table_y + Inches(0.4)
    for name, desc in rows:
        add_text(s, rx, y, col1_w, Inches(0.4),
                 name, size=10, color=INK, bold=True)
        add_text(s, rx + col1_w, y, col2_w, Inches(0.5),
                 desc, size=9.5, color=INK2, line_spacing=1.4)
        y += Inches(0.45)
        add_rule(s, rx, y - Inches(0.05), rw, color=LINE_SOFT, thickness=Pt(0.5))

    add_footer(s, "03  ·  PORTAL HOMES", "18 PÁGINAS  ·  ACCESO POR ROL")


def slide_capital_portal(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(8), Inches(0.6),
             "Portal Capital", font=SERIF, size=32, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "FINANCIAMIENTO PROPIO + INVERSORES", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.6), Inches(2.5), Inches(0.3),
             "/CAPITAL", font=MONO, size=8, color=INK4, letter_spacing=1.5, bold=True)
    add_text(s, MARGIN_X, Inches(1.95), Inches(2.5), Inches(1),
             "El dealer como banco",
             font=SERIF, size=15, color=GOLD, italic=True, line_spacing=1.3)

    rx = MARGIN_X + Inches(2.9)
    rw = CONTENT_W - Inches(2.9)
    add_text(s, rx, Inches(1.6), rw, Inches(1),
             "Cuando los bancos no financian — y para mobile homes pre-owned casi nunca lo hacen — el dealer se convierte en el prestamista. Capital es el portal donde se administra esa función. El acceso está restringido a un grupo reducido de personas autorizadas.",
             size=11, color=INK2, line_spacing=1.5)

    rows = [
        ("Dashboard", "Resumen de cartera, salud financiera, actividad reciente, KPIs."),
        ("Applications", "Solicitudes RTO en cada estado: pending, under review, approved, rejected, needs_info. Underwriting + scoring."),
        ("KYC", "Verificación manual de identidad. Documentos subidos por el cliente al bucket de Supabase, revisados aquí. Sin servicio externo."),
        ("Contracts", "Contratos RTO con sus 33 cláusulas, tabla de amortización, tracking de down payment por cuotas, generación de PDF."),
        ("Payments", "Pagos mensuales, confirmación de transferencias reportadas por el cliente, mora summary, comisiones, alertas de seguro vencido."),
        ("Investors", "Inversores activos, capital comprometido, capital desplegado, retornos."),
        ("Promissory notes", "Pagarés de inversores con alertas de maturity a 90, 60 y 30 días."),
        ("Accounting", "Contabilidad propia de Capital: chart of accounts, asientos, estados bancarios con reconciliación, reportes (P&L, balance, flujo de caja)."),
        ("Reports", "Estados mensuales en PDF para cada inversor + resúmenes unificados de cartera."),
        ("Mora", "Cartera vencida por aging bucket (0–30, 31–60, 61–90, 90+)."),
    ]
    table_y = Inches(2.65)
    col1_w = Inches(1.8)
    col2_w = rw - col1_w - Inches(0.2)
    add_text(s, rx, table_y, col1_w, Inches(0.25),
             "SECCIÓN", size=7, color=INK3, bold=True, letter_spacing=1.8)
    add_text(s, rx + col1_w, table_y, col2_w, Inches(0.25),
             "FUNCIÓN", size=7, color=INK3, bold=True, letter_spacing=1.8)
    add_rule(s, rx, table_y + Inches(0.28), rw)
    y = table_y + Inches(0.4)
    for name, desc in rows:
        add_text(s, rx, y, col1_w, Inches(0.36),
                 name, size=9.5, color=INK, bold=True)
        add_text(s, rx + col1_w, y, col2_w, Inches(0.5),
                 desc, size=9, color=INK2, line_spacing=1.4)
        y += Inches(0.36)
        add_rule(s, rx, y - Inches(0.04), rw, color=LINE_SOFT, thickness=Pt(0.5))

    add_footer(s, "04  ·  PORTAL CAPITAL", "19 PÁGINAS  ·  13 ARCHIVOS DE RUTAS BACKEND")


def slide_clientes_portal(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(8), Inches(0.6),
             "Portal Clientes", font=SERIF, size=32, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "CARA PÚBLICA Y AUTOSERVICIO", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.6), Inches(2.5), Inches(0.3),
             "/CLIENTES", font=MONO, size=8, color=INK4, letter_spacing=1.5, bold=True)
    add_text(s, MARGIN_X, Inches(1.95), Inches(2.5), Inches(1.2),
             "Del primer click al pago 36",
             font=SERIF, size=15, color=GOLD, italic=True, line_spacing=1.3)

    rx = MARGIN_X + Inches(2.9)
    rw = CONTENT_W - Inches(2.9)
    add_text(s, rx, Inches(1.6), rw, Inches(0.9),
             "Es el único portal accesible sin sesión. Un visitante puede llegar al catálogo desde Facebook o desde un link directo. La autenticación aparece solo al avanzar en el embudo de compra.",
             size=11, color=INK2, line_spacing=1.5)

    # pipeline (4 stages)
    add_text(s, rx, Inches(2.7), rw, Inches(0.3),
             "CAMINO DEL VISITANTE", size=8, color=GOLD, bold=True, letter_spacing=2)
    py = Inches(3.05)
    add_rule(s, rx, py, rw)
    add_rule(s, rx, py + Inches(1.4), rw)
    stages = [
        ("— 01 —", "Catálogo", "Listado público con filtros por ciudad y precio. Auto-refresh cada 2 min."),
        ("— 02 —", "Propiedad", "Detalle con galería, especificaciones y simulador RTO interactivo (down 30–100%, plazo 12–60 meses)."),
        ("— 03 —", "Compra", "Captura contacto, crea cliente o reusa existente, elige método: contado o RTO."),
        ("— 04 —", "Cuenta", "Estado de cuenta, KYC, firma de contrato, documentos, reportar pagos mensuales."),
    ]
    sw = rw / 4
    for i, (num, h, body) in enumerate(stages):
        x = rx + i * sw
        add_text(s, x + Inches(0.05), py + Inches(0.15), sw - Inches(0.2), Inches(0.25),
                 num, font=MONO, size=8, color=INK4, letter_spacing=1.5)
        add_text(s, x + Inches(0.05), py + Inches(0.4), sw - Inches(0.2), Inches(0.4),
                 h, font=SERIF, size=14, color=INK)
        add_text(s, x + Inches(0.05), py + Inches(0.78), sw - Inches(0.2), Inches(0.6),
                 body, size=8.5, color=INK3, line_spacing=1.4)
        if i < 3:
            add_vrule(s, x + sw, py, Inches(1.4))

    # bullet list below
    add_text(s, rx, Inches(4.7), rw, Inches(0.3),
             "LO QUE EL CLIENTE PUEDE HACER EN SU CUENTA",
             size=8, color=GOLD, bold=True, letter_spacing=2)
    items = [
        "Subir su licencia, pasaporte o ID estatal para KYC (JPG/PNG/WebP/HEIC, máx 10 MB)",
        "Firmar electrónicamente su contrato RTO con sello de tiempo, IP y user-agent registrados",
        "Completar la solicitud de crédito en secciones (empleo, vivienda, activos, deudas, referencias)",
        "Reportar el pago mensual del mes en curso y ver el historial completo de pagos",
        "Descargar contratos firmados y documentos de transferencia de título",
    ]
    y = Inches(5.1)
    for itm in items:
        add_text(s, rx, y, rw, Inches(0.3),
                 "—  " + itm, size=10, color=INK2, line_spacing=1.5)
        y += Inches(0.32)

    add_footer(s, "05  ·  PORTAL CLIENTES", "16 PÁGINAS  ·  PÚBLICO + AUTENTICADO")


def slide_loop(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(10), Inches(0.6),
             "Comprar · Renovar · Vender · Cobrar", font=SERIF, size=30, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "EL LOOP COMPLETO", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.6),
             "Cada etapa tiene su tabla, su(s) agente(s) de IA y su superficie de UI. Un mismo registro avanza por las cinco etapas sin salir del sistema y sin re-tipear datos.",
             size=11, color=INK2, line_spacing=1.5)

    # table
    rows = [
        ("1. Detectar", "BuscadorAgent recorre fuentes externas, aplica reglas de Maninos, propone candidatos.",
         "BuscadorAgent (Playwright + scrapers MHVillage, MobileHome.net, Zillow, FB Marketplace).",
         "market_listings"),
        ("2. Comprar", "Wizard “Revisar Casa”: fotos, checklist, datos del seller, oferta.",
         "FotosAgent (clasifica imágenes), VozAgent (notas dictadas en el yard).",
         "properties, title_transfers (purchase)"),
        ("3. Renovar", "Plan de renovación: materiales, mano de obra, contingencia.",
         "RenovacionAgent (orquestador), CostosAgent (precios desde DB).",
         "renovations · accounting_transactions"),
        ("4. Listar y vender", "Precio sugerido. Publicación. Cliente compra contado o solicita RTO.",
         "PrecioAgent (techo 80% del valor de mercado).",
         "sales, sale_payments, rto_applications"),
        ("5. Cobrar y cerrar", "Pagos mensuales, late fees, mora, transferencia final de título.",
         "Scheduler title_monitor diario contra TDHCA. Email service para recordatorios.",
         "rto_payments, title_transfers (sale), commission_payments"),
    ]
    ty = Inches(2.4)
    cols = [Inches(1.6), Inches(3.6), Inches(3.5), Inches(3.4)]
    headers = ["ETAPA", "QUÉ OCURRE", "AGENTES / SERVICIOS", "DATOS QUE PRODUCE"]
    cx = MARGIN_X
    for i, (col, hdr) in enumerate(zip(cols, headers)):
        add_text(s, cx, ty, col, Inches(0.25), hdr, size=7, color=INK3,
                 bold=True, letter_spacing=1.8)
        cx += col
    add_rule(s, MARGIN_X, ty + Inches(0.3), CONTENT_W)
    y = ty + Inches(0.4)
    for row in rows:
        cx = MARGIN_X
        for i, (col, val) in enumerate(zip(cols, row)):
            add_text(s, cx, y, col - Inches(0.15),
                     Inches(0.65),
                     val,
                     size=9 if i > 0 else 10,
                     color=INK if i == 0 else INK2,
                     bold=(i == 0),
                     font=MONO if i == 3 else SANS,
                     line_spacing=1.4)
            cx += col
        y += Inches(0.7)
        add_rule(s, MARGIN_X, y - Inches(0.05), CONTENT_W,
                 color=LINE_SOFT, thickness=Pt(0.5))

    add_footer(s, "06  ·  EL LOOP", "5 ETAPAS  ·  6 AGENTES  ·  UNA FUENTE DE VERDAD")


def slide_agents(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(10), Inches(0.6),
             "Los seis agentes especializados", font=SERIF, size=30, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "CAPA DE IA", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.45), CONTENT_W, Inches(0.7),
             "Cada agente vive en api/agents/, hereda de BaseAgent, opera de forma asíncrona y devuelve JSON estructurado validado por Pydantic. No tienen memoria entre llamadas: el contexto necesario se pasa explícitamente en cada request.",
             size=11, color=INK2, line_spacing=1.5)

    rows = [
        ("Buscador", "Encuentra mobile homes en el mercado que cumplen las reglas de Maninos.",
         "Scrapea MHVillage, MobileHome.net, MHBay, Zillow, Facebook Marketplace. Filtra por reglas Feb 2026: 60% del valor de mercado, $5K–$80K, 200mi de Houston o Dallas, sin filtro de año."),
        ("Costos", "Estima el costo de la renovación.",
         "Lee el catálogo de materiales desde la base de datos. Devuelve materiales + mano de obra (30–50% material) + contingencia (5–15%) + costo por sqft + factores de riesgo."),
        ("Precio", "Propone precio de venta.",
         "Devuelve cuatro estrategias: mínimo, target, mercado, agresivo. Aplica techo automático del 80% del valor de mercado."),
        ("Fotos", "Clasifica fotografías por estado de renovación.",
         "Vision real (no keywords). Distingue before / after / exterior / interior. Ordena la galería pública: after primero, después exterior, interior, before."),
        ("Voz", "Transcribe audio del staff y extrae intención.",
         "Whisper en español. Detecta si es una nota de inspección, una lista de materiales o una observación. Funciona con Spanglish."),
        ("Renovación", "Conversación guiada para planificar la renovación.",
         "Tiene los precios de los materiales embebidos en el prompt (decisión de latencia). Si no hay sqft, no estima — falla loud."),
    ]
    ty = Inches(2.45)
    cols = [Inches(1.4), Inches(3.5), Inches(7.2)]
    headers = ["AGENTE", "FUNCIÓN", "DETALLE RELEVANTE"]
    cx = MARGIN_X
    for col, hdr in zip(cols, headers):
        add_text(s, cx, ty, col, Inches(0.25), hdr, size=7, color=INK3,
                 bold=True, letter_spacing=1.8)
        cx += col
    add_rule(s, MARGIN_X, ty + Inches(0.3), CONTENT_W)
    y = ty + Inches(0.4)
    for row in rows:
        cx = MARGIN_X
        for i, (col, val) in enumerate(zip(cols, row)):
            add_text(s, cx, y, col - Inches(0.15), Inches(0.55),
                     val,
                     size=9 if i > 0 else 10,
                     color=INK if i == 0 else INK2,
                     bold=(i == 0),
                     line_spacing=1.4)
            cx += col
        y += Inches(0.55)
        add_rule(s, MARGIN_X, y - Inches(0.04), CONTENT_W,
                 color=LINE_SOFT, thickness=Pt(0.5))

    # divider
    add_rule(s, MARGIN_X, Inches(6.4), CONTENT_W)
    add_text(s, MARGIN_X, Inches(6.55), CONTENT_W, Inches(0.7),
             "Aparte: el chat flotante AIChatWidget que aparece en el portal Homes no usa estos agentes. Llama a /api/ai/chat (gpt-5-mini) con un patrón de tool-calling sobre 19 herramientas de consulta a la base de datos.",
             size=10, color=INK2, italic=True, line_spacing=1.5)

    add_footer(s, "07  ·  AGENTES", "6 ESPECIALIZADOS  +  1 ASISTENTE CONVERSACIONAL")


def slide_scheduler(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(8), Inches(0.6),
             "El scheduler", font=SERIF, size=32, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "AUTOMATIZACIÓN EN BACKGROUND", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.7),
             "APScheduler corriendo en zona horaria US Central. Diez jobs en producción. Cada ejecución, exitosa o fallida, deja registro en la tabla scheduler_runs — el widget en /homes/transfers lee de ahí, no de la memoria del scheduler.",
             size=11, color=INK2, line_spacing=1.5)

    rows = [
        ("process_scheduled_emails", "Cada 30 min", "Envía emails programados pendientes."),
        ("rto_reminders", "Diario · 8:00 CT", "Recordatorios de pago RTO al cliente (3d antes, día del pago, 1d después)."),
        ("rto_overdue_alerts", "Diario · 9:00 CT", "Alertas internas para pagos RTO vencidos."),
        ("promissory_maturity_alerts", "Diario · 9:30 CT", "Avisa al equipo de pagarés a vencer (90 / 60 / 30 días)."),
        ("title_monitor", "Diario · 10:00 CT", "Consulta TDHCA por cada serial pendiente, fuzzy-match del nombre."),
        ("investor_followup_emails", "Día 1 · 10:30 CT", "Estado mensual a cada inversor con su PDF."),
        ("portal_sync", "Cada 2 horas", "Consistencia cruzada Homes ↔ Capital."),
        ("refresh_partner_listings", "Cada 6 horas", "Refresca listings de VMF Homes y 21st Mortgage."),
        ("facebook_auto_scrape", "Lun y Jue · 7:00 CT", "Scraping de Facebook Marketplace vía Apify."),
        ("expire_old_listings", "Diario · 6:00 CT", "Marca como expirados listings con más de 14 días."),
    ]
    ty = Inches(2.45)
    cols = [Inches(3.7), Inches(2.3), Inches(6.1)]
    headers = ["JOB", "CADENCIA", "FUNCIÓN"]
    cx = MARGIN_X
    for col, hdr in zip(cols, headers):
        add_text(s, cx, ty, col, Inches(0.25), hdr, size=7, color=INK3,
                 bold=True, letter_spacing=1.8)
        cx += col
    add_rule(s, MARGIN_X, ty + Inches(0.3), CONTENT_W)
    y = ty + Inches(0.38)
    for row in rows:
        cx = MARGIN_X
        for i, (col, val) in enumerate(zip(cols, row)):
            add_text(s, cx, y, col - Inches(0.15), Inches(0.4),
                     val,
                     size=9.5 if i != 0 else 10,
                     color=GOLD_SOFT if i == 0 else INK2,
                     bold=(i == 0),
                     font=MONO if i == 0 else SANS,
                     line_spacing=1.3)
            cx += col
        y += Inches(0.36)
        add_rule(s, MARGIN_X, y - Inches(0.03), CONTENT_W,
                 color=LINE_SOFT, thickness=Pt(0.5))

    add_footer(s, "08  ·  SCHEDULER", "10 JOBS  ·  US CENTRAL  ·  AUDIT EN scheduler_runs")


def slide_titles(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(10), Inches(0.6),
             "El flujo del título estatal", font=SERIF, size=30, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "TDHCA Y HERENCIA", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.7),
             "En Texas, la transferencia de propiedad de un mobile home se registra en TDHCA mediante el Statement of Ownership. Maninos OS automatiza el seguimiento.",
             size=11, color=INK2, line_spacing=1.5)

    col_w = (CONTENT_W - Inches(0.6)) / 2
    col1_x = MARGIN_X
    col2_x = MARGIN_X + col_w + Inches(0.6)

    add_text(s, col1_x, Inches(2.5), col_w, Inches(0.3),
             "FLUJO AUTOMÁTICO", size=8, color=GOLD, bold=True, letter_spacing=2)
    auto_steps = [
        "1.  Al subir el documento de aplicación de título, el agente extrae el HUD label / serial.",
        "2.  Si la aplicación de título no tiene el dato, fallback al Bill of Sale (bos_purchase.hud_label_number).",
        "3.  El job diario de las 10:00 CT consulta el portal TDHCA por cada serial pendiente.",
        "4.  Compara el nombre que devuelve TDHCA contra to_name con un fuzzy-match.",
        "5.  Si coinciden, marca title_name_updated = TRUE, registra el nombre en tdhca_owner_name, mueve la transferencia a completed.",
    ]
    y = Inches(2.85)
    for stp in auto_steps:
        add_text(s, col1_x, y, col_w, Inches(0.6), stp,
                 size=10, color=INK2, line_spacing=1.5)
        y += Inches(0.55)

    add_text(s, col2_x, Inches(2.5), col_w, Inches(0.3),
             "CASAS LEGACY", size=8, color=GOLD, bold=True, letter_spacing=2)
    add_text(s, col2_x, Inches(2.85), col_w, Inches(0.7),
             "Las casas vendidas antes de la plataforma no están en el sistema y su transferencia ya se hizo fuera. Para incorporarlas, el portal Homes ofrece la subida manual:",
             size=10, color=INK2, line_spacing=1.5)
    legacy_items = [
        "Permite crear la propiedad y la transferencia desde cero o desde una propiedad existente.",
        "Si la casa está vendida, pide el comprador y crea la transferencia de venta correspondiente.",
        "Marca is_manual_upload = TRUE y title_name_updated = TRUE para que el job diario las ignore.",
        "El serial mostrado en el listado es un link directo a la ficha pública de TDHCA.",
    ]
    y = Inches(3.85)
    for itm in legacy_items:
        add_text(s, col2_x, y, col_w, Inches(0.5),
                 "—  " + itm, size=10, color=INK2, line_spacing=1.5)
        y += Inches(0.45)

    add_footer(s, "09  ·  TÍTULOS", "TDHCA + MANUAL  ·  AUDIT CADA NOCHE")


def slide_rto(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(10), Inches(0.6),
             "El módulo de financiamiento propio", font=SERIF, size=30, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "RTO EN DETALLE", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.7),
             "El contrato base es el Texas Residential Lease With Purchase Option (33 cláusulas). Los defaults numéricos están escritos en el código y son configurables por contrato.",
             size=11, color=INK2, line_spacing=1.5)

    # KPI row
    add_rule(s, MARGIN_X, Inches(2.5), CONTENT_W)
    add_rule(s, MARGIN_X, Inches(3.4), CONTENT_W)
    kpis = [
        ("$15", "/día", "LATE FEE"),
        ("5", " días", "PERÍODO DE GRACIA"),
        ("15", "", "DÍA DEL MES DE VENCIMIENTO"),
        ("$250", "", "NSF FEE"),
        ("$695", "/mes", "HOLDOVER"),
    ]
    kw = CONTENT_W / 5
    for i, (num, suffix, lbl) in enumerate(kpis):
        x = MARGIN_X + i * kw
        add_text(s, x + Inches(0.05), Inches(2.65), kw, Inches(0.5),
                 [(num, {"font": SERIF, "size": 24, "color": INK}),
                  (suffix, {"font": SANS, "size": 11, "color": INK3})],
                 line_spacing=1.0)
        add_text(s, x + Inches(0.05), Inches(3.2), kw, Inches(0.2),
                 lbl, size=7, color=INK3, bold=True, letter_spacing=1.8)
        if i < 4:
            add_vrule(s, x + kw, Inches(2.55), Inches(0.85))

    add_text(s, MARGIN_X, Inches(3.55), CONTENT_W, Inches(0.3),
             "Defaults aplicados por la migración 012_maninos_capital.sql. Cada contrato puede sobrescribirlos.",
             size=9, color=INK3, italic=True)

    add_rule(s, MARGIN_X, Inches(4.05), CONTENT_W)

    col_w = (CONTENT_W - Inches(0.6)) / 2
    col1_x = MARGIN_X
    col2_x = MARGIN_X + col_w + Inches(0.6)

    add_text(s, col1_x, Inches(4.25), col_w, Inches(0.3),
             "LO QUE EL DEALER VE CADA DÍA",
             size=8, color=GOLD, bold=True, letter_spacing=2)
    dealer_items = [
        "Solicitudes RTO entrantes con su scoring de underwriting",
        "Pagos del mes con su confirmación pendiente o aplicada",
        "Cartera vencida segmentada por aging (0–30, 31–60, 61–90, 90+)",
        "Comisiones por venta dividas finder vs closer, pendientes vs pagadas",
        "Alertas de seguros del cliente próximos a vencer",
    ]
    y = Inches(4.6)
    for itm in dealer_items:
        add_text(s, col1_x, y, col_w, Inches(0.4),
                 "—  " + itm, size=10, color=INK2, line_spacing=1.5)
        y += Inches(0.35)

    add_text(s, col2_x, Inches(4.25), col_w, Inches(0.3),
             "LO QUE EL CLIENTE VE EN SU CUENTA",
             size=8, color=GOLD, bold=True, letter_spacing=2)
    client_items = [
        "Su tabla de amortización con meses pasados, presente y futuros",
        "Total pagado y saldo restante",
        "Botón para reportar el pago del mes con método y referencia",
        "Estado de su KYC y su contrato",
        "Notificaciones de pagos próximos y vencidos",
    ]
    y = Inches(4.6)
    for itm in client_items:
        add_text(s, col2_x, y, col_w, Inches(0.4),
                 "—  " + itm, size=10, color=INK2, line_spacing=1.5)
        y += Inches(0.35)

    add_footer(s, "10  ·  RTO", "33 CLÁUSULAS  ·  DOBLE VISTA (OPERADOR + CLIENTE)")


def slide_accounting(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(8), Inches(0.6),
             "La capa contable", font=SERIF, size=32, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "GL COMPLETO, DOS ENTIDADES", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.9),
             "Maninos OS no se apoya en QuickBooks externo: implementa su propio General Ledger con una estructura inspirada en el chart of accounts de QuickBooks. Hay dos contabilidades separadas — una para Homes, otra para Capital — porque son dos LLCs distintas con reportes independientes.",
             size=11, color=INK2, line_spacing=1.5)

    col_w = (CONTENT_W - Inches(0.6)) / 2
    col1_x = MARGIN_X
    col2_x = MARGIN_X + col_w + Inches(0.6)

    add_text(s, col1_x, Inches(2.7), col_w, Inches(0.3),
             "TABLAS IMPLICADAS", size=8, color=GOLD, bold=True, letter_spacing=2)
    schema_lines = [
        ("accounting_accounts",       "// chart of accounts Homes"),
        ("accounting_transactions",   "// asientos contables"),
        ("bank_accounts",             "// cuentas bancarias"),
        ("recurring_expenses",        "// reglas recurrentes"),
        ("accounting_budgets",        "// presupuestos"),
        ("accounting_audit_log",      "// auditoría"),
        ("capital_accounts",          "// chart of accounts Capital"),
        ("capital_transactions",      "// asientos Capital"),
        ("capital_bank_statements",   "// extractos bancarios"),
        ("capital_statement_movements","// líneas de extracto"),
        ("capital_budgets",           "// presupuestos Capital"),
    ]
    # bg block for schema
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, col1_x, Inches(3.05), col_w, Inches(3.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = SURFACE
    bg.line.color.rgb = LINE
    bg.line.width = Pt(0.5)
    bg.shadow.inherit = False

    y = Inches(3.2)
    for key, comment in schema_lines:
        add_text(s, col1_x + Inches(0.2), y, col_w - Inches(0.4), Inches(0.3),
                 [(key.ljust(30), {"font": MONO, "size": 9, "color": GOLD_SOFT}),
                  (comment, {"font": MONO, "size": 9, "color": INK4})],
                 line_spacing=1.0)
        y += Inches(0.3)

    add_text(s, col2_x, Inches(2.7), col_w, Inches(0.3),
             "CATEGORÍAS QUE CUBRE",
             size=8, color=GOLD, bold=True, letter_spacing=2)
    cat_items = [
        "Ingresos: ventas contado, ventas RTO, pagos mensuales",
        "Gastos: compra, renovación, transporte, comisiones, operativos",
        "Activos: inventario por casa (cada propiedad genera su jerarquía)",
        "Reportes: P&L, balance, flujo de caja, presupuesto vs real, P&L por propiedad y por yard",
        "Reconciliación bancaria con import de extractos y matching",
        "Audit log para cumplimiento y trazabilidad",
    ]
    y = Inches(3.05)
    for itm in cat_items:
        add_text(s, col2_x, y, col_w, Inches(0.5),
                 "—  " + itm, size=10, color=INK2, line_spacing=1.5)
        y += Inches(0.4)

    add_text(s, col2_x, Inches(5.95), col_w, Inches(0.5),
             "Los pagos RTO disparan automáticamente el asiento contable correspondiente (_accounting_hooks.py). El equipo no rehace contabilidad.",
             size=10, color=INK2, italic=True, line_spacing=1.5)

    add_footer(s, "11  ·  CONTABILIDAD", "DOS LLCs  ·  GL PROPIO  ·  AUDIT LOG")


def slide_stack(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(8), Inches(0.6),
             "El stack técnico", font=SERIF, size=32, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "SIN SORPRESAS, TODO ESTÁNDAR", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    rows = [
        ("Backend", "Python 3.12 con FastAPI. 22 archivos de rutas en Homes, 13 en Capital, 5 públicos. 9 servicios (email, pdf, scheduler, scrapers, title monitor, notification, document, esign, property)."),
        ("Datos", "Supabase: Postgres + Auth + Storage. 58 tablas, 91 migraciones SQL. Acceso vía SDK directo, sin ORM. Service-role key en backend, RLS en buckets públicos."),
        ("Frontend", "Next.js 14 con App Router. Tailwind con sistema de diseño propio (paleta navy / oro). Mobile-first responsive. Sin librería de estado externa: hooks de React."),
        ("IA", "OpenAI: gpt-5 en agentes, gpt-5-mini en chat con tool-calling, whisper-1 para transcripción (forzado a español)."),
        ("Email", "Resend con templates propios y attachments PDF. Procesamiento de cola de emails programados cada 30 min."),
        ("Scraping", "Playwright para automatización de browser. Apify para Facebook Marketplace."),
        ("Despliegue", "Backend en Railway (Docker). Frontend en Vercel. Deploy automático en cada push a main."),
        ("Observabilidad", "structlog para logging estructurado. Logfire opcional para telemetría."),
    ]
    ty = Inches(1.6)
    cols = [Inches(2.0), Inches(10.0)]
    headers = ["CAPA", "COMPONENTES"]
    cx = MARGIN_X
    for col, hdr in zip(cols, headers):
        add_text(s, cx, ty, col, Inches(0.25), hdr, size=7, color=INK3,
                 bold=True, letter_spacing=1.8)
        cx += col
    add_rule(s, MARGIN_X, ty + Inches(0.3), CONTENT_W)
    y = ty + Inches(0.4)
    for capa, comp in rows:
        add_text(s, MARGIN_X, y, cols[0] - Inches(0.15), Inches(0.55),
                 capa, size=11, color=INK, bold=True)
        add_text(s, MARGIN_X + cols[0], y, cols[1] - Inches(0.15), Inches(0.7),
                 comp, size=10, color=INK2, line_spacing=1.45)
        y += Inches(0.6)
        add_rule(s, MARGIN_X, y - Inches(0.04), CONTENT_W,
                 color=LINE_SOFT, thickness=Pt(0.5))

    add_footer(s, "12  ·  STACK", "FASTAPI + NEXT.JS + SUPABASE")


def slide_rules(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(10), Inches(0.6),
             "Las reglas del negocio", font=SERIF, size=32, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "TEXAS, MANINOS, 2026", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.7),
             "Las reglas viven en el código y se aplican automáticamente. No son recomendaciones del agente: son guardarraíles que las APIs y los agentes respetan.",
             size=11, color=INK2, line_spacing=1.5)

    col_w = (CONTENT_W - Inches(0.6)) / 2
    col1_x = MARGIN_X
    col2_x = MARGIN_X + col_w + Inches(0.6)

    add_text(s, col1_x, Inches(2.7), col_w, Inches(0.3),
             "ADQUISICIÓN", size=8, color=GOLD, bold=True, letter_spacing=2)
    adq_rules = [
        ("Geografía:", "radio de 200 millas desde Houston o Dallas, solo Texas"),
        ("Precio de compra:", "hasta el 60% del valor de mercado (renovación no incluida)"),
        ("Rango aceptado:", "entre $5,000 y $80,000 por propiedad"),
        ("Año:", "sin filtro"),
        ("Tipos:", "single wide y double wide"),
    ]
    y = Inches(3.05)
    for k, v in adq_rules:
        add_text(s, col1_x, y, col_w, Inches(0.4),
                 [("—  ", {"color": GOLD, "size": 11}),
                  (k + " ", {"size": 11, "color": INK, "bold": True}),
                  (v, {"size": 11, "color": INK2})],
                 line_spacing=1.5)
        y += Inches(0.45)

    add_text(s, col2_x, Inches(2.7), col_w, Inches(0.3),
             "RENOVACIÓN Y VENTA", size=8, color=GOLD, bold=True, letter_spacing=2)
    sell_rules = [
        ("Presupuesto de renovación:", "entre $5,000 y $15,000"),
        ("Precio de venta:", "hasta el 80% del valor de mercado mezclado (techo enforced por PrecioAgent)"),
        ("Comisiones:", "$1,500 en venta contado · $1,000 en venta RTO"),
        ("Reparto de comisión:", "50% para el finder, 50% para el closer"),
        ("Modalidades de venta:", "contado o RTO"),
    ]
    y = Inches(3.05)
    for k, v in sell_rules:
        add_text(s, col2_x, y, col_w, Inches(0.5),
                 [("—  ", {"color": GOLD, "size": 11}),
                  (k + " ", {"size": 11, "color": INK, "bold": True}),
                  (v, {"size": 11, "color": INK2})],
                 line_spacing=1.5)
        y += Inches(0.45)

    add_footer(s, "13  ·  REGLAS", "APLICADAS EN CÓDIGO  ·  FEBRERO 2026")


def slide_portability(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(11), Inches(0.6),
             "Lo que es genérico, lo que es Texas-MH", font=SERIF, size=28, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "ANÁLISIS DE PORTABILIDAD", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.9),
             "El núcleo del sistema describe un patrón universal: un dealer adquiere un activo físico de alto valor con título registrable, lo reacondiciona, lo vende a crédito propio, y cobra en el tiempo. Cualquier negocio que cumpla ese patrón puede correr sobre Maninos OS con configuración localizada.",
             size=11, color=INK2, italic=True, line_spacing=1.5)

    col_w = (CONTENT_W - Inches(0.6)) / 2
    col1_x = MARGIN_X
    col2_x = MARGIN_X + col_w + Inches(0.6)

    add_text(s, col1_x, Inches(2.95), col_w, Inches(0.3),
             "GENÉRICO — EL MOTOR", size=8, color=GOLD, bold=True, letter_spacing=2)
    gen_items = [
        "Pipeline de adquisición: detectar, calificar, comprar",
        "Workflow de reacondicionamiento con materiales, mano de obra y costos",
        "Galería de fotos clasificada (vision agnóstico al activo)",
        "Motor RTO/BHPH: contrato, amortización, late fee, NSF, holdover, mora",
        "Portal del cliente: KYC, firma electrónica, reporte de pagos, estado de cuenta",
        "Capital de inversores con promissory notes y reportes mensuales",
        "General Ledger con reconciliación bancaria",
        "Asistente conversacional sobre la base de datos",
    ]
    y = Inches(3.3)
    for itm in gen_items:
        add_text(s, col1_x, y, col_w, Inches(0.4),
                 "—  " + itm, size=10, color=INK2, line_spacing=1.5)
        y += Inches(0.38)

    add_text(s, col2_x, Inches(2.95), col_w, Inches(0.3),
             "LOCALIZABLE — LA «ASSET PACK»", size=8, color=GOLD, bold=True, letter_spacing=2)
    loc_items = [
        ("Schema del activo:", "sqft + HUD label en MH; HIN + nº motor en boats; VIN + odómetro en autos."),
        ("Reglas de adquisición:", "el 60% / $5K-$80K / 200mi es de Maninos; otro dealer las suyas."),
        ("Workflow de título:", "TDHCA en Texas, otra agencia en cada estado, USCG opcional para boats grandes."),
        ("Plantilla de contrato:", "Lease with Purchase Option en MH; contrato BHPH en autos."),
        ("Fuentes de scraping:", "MHVillage para MH, Manheim para autos, BoatTrader para boats."),
        ("Prompts de los agentes:", "agente de costos entiende pintura y tablaroca; gelcoat y bottom paint en boats."),
    ]
    y = Inches(3.3)
    for k, v in loc_items:
        add_text(s, col2_x, y, col_w, Inches(0.5),
                 [("—  ", {"color": GOLD, "size": 10}),
                  (k + " ", {"size": 10, "color": INK, "bold": True}),
                  (v, {"size": 10, "color": INK2})],
                 line_spacing=1.5)
        y += Inches(0.42)

    add_footer(s, "14  ·  PORTABILIDAD", "NÚCLEO + ASSET PACK")


def slide_verticals(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(11), Inches(0.6),
             "Verticales donde el patrón calza", font=SERIF, size=30, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "LECTURA CUALITATIVA", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.7),
             "No se incluyen tamaños de mercado en este documento porque son estimaciones de terceros que no he verificado. Lo que sigue es una lectura cualitativa basada en cuán literal es la analogía con el caso Maninos.",
             size=10, color=INK3, italic=True, line_spacing=1.5)

    rows = [
        ("BHPH used cars",
         "El loop comprar → reacondicionar → financiar → cobrar es idéntico. La cartera vencida y el tracking de pagos mensuales son la analogía más directa.",
         "Workflow de título por estado vía DMV. Integraciones starter-interrupt / GPS. Auctions como fuente."),
        ("RV / travel trailers",
         "El concepto de “Revisar Casa” se mapea casi sin cambios al PDI (Pre-Delivery Inspection). El ciclo es el mismo.",
         "Schema del activo (motorhome vs towable). Plantilla de inspección PDI."),
        ("Boats (used)",
         "Reacondicionamiento estándar (gelcoat, motor, tapizado). Título estatal análogo. Financiamiento propio es el vacío más grande.",
         "HIN + datos de motor / trailer en el schema. Estacionalidad. USCG para vessels > 5 net tons."),
        ("Tiny homes / park models / ADU",
         "Casi 1:1 con MH. Mismo título en muchos estados.",
         "Volumen pequeño; tiene más sentido como módulo dentro de un dealer MH que como vertical independiente."),
        ("Maquinaria pesada / tractores",
         "Adquisición y reacondicionamiento encajan. El financiamiento propio es raro porque la captiva del fabricante domina.",
         "UCC-1 en lugar de título. Lock-in de OEM hace el GTM más difícil."),
    ]
    ty = Inches(2.65)
    cols = [Inches(2.6), Inches(5.0), Inches(4.5)]
    headers = ["VERTICAL", "CÓMO ENCAJA CON EL MOTOR ACTUAL", "QUÉ HAY QUE ADAPTAR"]
    cx = MARGIN_X
    for col, hdr in zip(cols, headers):
        add_text(s, cx, ty, col, Inches(0.25), hdr, size=7, color=INK3,
                 bold=True, letter_spacing=1.8)
        cx += col
    add_rule(s, MARGIN_X, ty + Inches(0.3), CONTENT_W)
    y = ty + Inches(0.4)
    for row in rows:
        cx = MARGIN_X
        for i, (col, val) in enumerate(zip(cols, row)):
            add_text(s, cx, y, col - Inches(0.15), Inches(0.7),
                     val,
                     size=9.5 if i > 0 else 11,
                     color=INK if i == 0 else INK2,
                     bold=(i == 0),
                     line_spacing=1.4)
            cx += col
        y += Inches(0.7)
        add_rule(s, MARGIN_X, y - Inches(0.04), CONTENT_W,
                 color=LINE_SOFT, thickness=Pt(0.5))

    add_footer(s, "15  ·  VERTICALES", "LECTURA CUALITATIVA  ·  SIN CIFRAS DE TERCEROS")


def slide_status(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(10), Inches(0.6),
             "Dónde está hoy Maninos OS", font=SERIF, size=30, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "ESTADO VERIFICABLE", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    add_text(s, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.7),
             "Lo siguiente es lo que hay en el repositorio y en producción ahora mismo, contado contra el código actual.",
             size=11, color=INK2, line_spacing=1.5)

    add_rule(s, MARGIN_X, Inches(2.4), CONTENT_W)
    add_rule(s, MARGIN_X, Inches(3.3), CONTENT_W)
    kpis = [
        ("3", "PORTALES EN PRODUCCIÓN"),
        ("53", "PÁGINAS FRONTEND (18+19+16)"),
        ("40", "ARCHIVOS DE RUTAS BACKEND"),
        ("6", "AGENTES ESPECIALIZADOS"),
        ("10", "JOBS EN SCHEDULER"),
    ]
    kw = CONTENT_W / 5
    for i, (num, lbl) in enumerate(kpis):
        x = MARGIN_X + i * kw
        add_text(s, x + Inches(0.05), Inches(2.55), kw, Inches(0.5),
                 num, font=SERIF, size=28, color=INK)
        add_text(s, x + Inches(0.05), Inches(3.05), kw, Inches(0.2),
                 lbl, size=7, color=INK3, bold=True, letter_spacing=1.8)
        if i < 4:
            add_vrule(s, x + kw, Inches(2.45), Inches(0.85))

    add_rule(s, MARGIN_X, Inches(3.7), CONTENT_W)

    col_w = (CONTENT_W - Inches(0.6)) / 2
    col1_x = MARGIN_X
    col2_x = MARGIN_X + col_w + Inches(0.6)

    add_text(s, col1_x, Inches(3.95), col_w, Inches(0.3),
             "EN PRODUCCIÓN", size=8, color=GOLD, bold=True, letter_spacing=2)
    prod_items = [
        "Operación real diaria del staff de Maninos en el portal Homes",
        "Inversores activos consultando sus reportes en Capital",
        "Clientes reportando pagos mensuales desde su cuenta en Clientes",
        "Scheduler diario contra TDHCA dejando audit en scheduler_runs",
        "Agentes de IA en flujos de “Revisar Casa” y planificación de renovación",
    ]
    y = Inches(4.3)
    for itm in prod_items:
        add_text(s, col1_x, y, col_w, Inches(0.4),
                 "—  " + itm, size=10, color=INK2, line_spacing=1.5)
        y += Inches(0.4)

    add_text(s, col2_x, Inches(3.95), col_w, Inches(0.3),
             "ÁREAS QUE AÚN TIENEN DEUDA",
             size=8, color=GOLD, bold=True, letter_spacing=2)
    debt_items = [
        "Multi-tenant: hoy es instancia única para Maninos. Falta separación de tenants.",
        "White-label de marca y dominio.",
        "Stripe está como label en dropdowns pero no integrado de verdad. Pagos online aún no existen — todo es transferencia con confirmación.",
        "Asset Pack como abstracción explícita: hoy las reglas MH están dispersas en código; portarlas requiere refactor.",
    ]
    y = Inches(4.3)
    for itm in debt_items:
        add_text(s, col2_x, y, col_w, Inches(0.5),
                 "—  " + itm, size=10, color=INK2, line_spacing=1.5)
        y += Inches(0.5)

    add_footer(s, "16  ·  ESTADO ACTUAL", "PRODUCCIÓN REAL  ·  DEUDA EXPLÍCITA")


def slide_closing(prs):
    s = slide_blank(prs)
    add_text(s, MARGIN_X, MARGIN_Y, Inches(8), Inches(0.6),
             "En una frase", font=SERIF, size=32, color=INK)
    add_text(s, MARGIN_X, MARGIN_Y, CONTENT_W, Inches(0.6),
             "CIERRE", size=8, color=INK3, bold=True,
             letter_spacing=2.5, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, MARGIN_X, MARGIN_Y + Inches(0.7), CONTENT_W)

    # vertical gold rule + pulled quote
    quote_x = MARGIN_X + Inches(0.4)
    quote_y = Inches(2.0)
    quote_w = CONTENT_W - Inches(0.4)
    add_vrule(s, quote_x - Inches(0.1), quote_y, Inches(2.6), color=GOLD, thickness=Pt(2.5))
    add_text(s, quote_x + Inches(0.2), quote_y, quote_w - Inches(0.2), Inches(2.6),
             "“Maninos OS es la disciplina operativa de un dealer especializado de mobile homes en Texas — convertida en código. Lo que era spreadsheet, WhatsApp, QuickBooks y procedimientos no escritos, ahora es un sistema con seis agentes de IA, GL propio, módulo de financiamiento, portal del cliente y audit log.”",
             font=SERIF, size=21, color=INK, italic=True, line_spacing=1.4)

    add_rule(s, MARGIN_X, Inches(4.95), CONTENT_W)

    add_text(s, MARGIN_X, Inches(5.2), CONTENT_W, Inches(0.9),
             "El siguiente paso natural es separar el motor de las reglas específicas de Maninos, formalizar el concepto de Asset Pack, y permitir que otros dealers — primero de mobile homes en otros estados, después de RV o BHPH — corran sobre la misma infraestructura.",
             font=SERIF, size=15, color=INK2, line_spacing=1.5)

    # bottom
    add_rule(s, MARGIN_X, Inches(6.6), CONTENT_W)
    add_text(s, MARGIN_X, Inches(6.8), Inches(8), Inches(0.4),
             "Maninos OS", font=SERIF, size=16, color=INK)
    add_text(s, MARGIN_X, Inches(7.05), Inches(8), Inches(0.3),
             "Documento descriptivo · 2026", size=9, color=INK3, italic=True)
    add_text(s, MARGIN_X, Inches(6.8), CONTENT_W, Inches(0.3),
             "Maninos Homes LLC", size=9, color=INK3,
             align=PP_ALIGN.RIGHT)
    add_text(s, MARGIN_X, Inches(7.05), CONTENT_W, Inches(0.3),
             "Maninos Capital LLC", size=9, color=INK3,
             align=PP_ALIGN.RIGHT)

    add_footer(s, "17  ·  CIERRE", "FIN DEL DOCUMENTO")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_definicion(prs)
    slide_arquitectura(prs)
    slide_section_index(prs, "— PARTE I —", "El sistema, módulo por módulo",
                        "Recorrido descriptivo por cada portal: qué incluye, cómo lo usa cada usuario, qué datos guarda y cómo se conecta al resto.")
    slide_homes_portal(prs)
    slide_capital_portal(prs)
    slide_clientes_portal(prs)
    slide_section_index(prs, "— PARTE II —", "El loop de negocio",
                        "Cómo se mueve un activo desde que aparece en una fuente externa hasta que termina su contrato RTO 36 meses después.")
    slide_loop(prs)
    slide_agents(prs)
    slide_scheduler(prs)
    slide_titles(prs)
    slide_rto(prs)
    slide_accounting(prs)
    slide_stack(prs)
    slide_rules(prs)
    slide_section_index(prs, "— PARTE III —", "El patrón es portable",
                        "El sistema fue construido para mobile homes en Texas, pero su modelo de datos y sus flujos describen una clase más amplia de negocios.")
    slide_portability(prs)
    slide_verticals(prs)
    slide_status(prs)
    slide_closing(prs)

    out = os.path.join(os.path.dirname(__file__), "Maninos_OS_Pitch.pptx")
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
