#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Maninos AI
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Colors
NAVY_900 = RGBColor(15, 23, 42)
NAVY_800 = RGBColor(30, 41, 59)
NAVY_700 = RGBColor(51, 65, 85)
NAVY_100 = RGBColor(241, 245, 249)
GOLD_500 = RGBColor(212, 168, 83)
GOLD_400 = RGBColor(229, 190, 115)
GREEN_500 = RGBColor(34, 197, 94)
WHITE = RGBColor(255, 255, 255)

def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_text(slide, text, left, top, width, height, font_size=44, bold=True, color=WHITE):
    """Add a text box with title styling"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return txBox

def add_bullet_list(slide, items, left, top, width, height, font_size=16, color=NAVY_100):
    """Add a bulleted list"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    
    return txBox

def add_card(slide, left, top, width, height, title, items, accent_color=GOLD_500):
    """Add a card with title and bullet points"""
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_800
    shape.line.color.rgb = NAVY_700
    
    # Title
    add_title_text(slide, title, left + 0.2, top + 0.15, width - 0.4, 0.4, font_size=14, bold=True, color=accent_color)
    
    # Items
    add_bullet_list(slide, items, left + 0.2, top + 0.5, width - 0.4, height - 0.6, font_size=12, color=NAVY_100)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY_900)
    
    # Logo text
    add_title_text(slide, "MANINOS CAPITAL LLC", 0.5, 0.5, 5, 0.5, font_size=14, bold=False, color=GOLD_500)
    
    # Main title
    add_title_text(slide, "Maninos AI", 0.5, 2.5, 12, 1.2, font_size=72, bold=True, color=WHITE)
    
    # Subtitle
    add_title_text(slide, "Plataforma Inteligente para Rent-to-Own", 0.5, 3.8, 12, 0.6, font_size=28, bold=False, color=NAVY_700)
    
    # Date badge
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(5.2), Inches(4.3), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD_500
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = "📅 Enero 2026 - Desarrollo Activo"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY_900
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # =========================================================================
    # SLIDE 2: Progress Overview
    # =========================================================================
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY_900)
    
    add_title_text(slide, "📊 Progreso del Desarrollo", 0.5, 0.3, 12, 0.8, font_size=36, bold=True, color=GOLD_500)
    
    # Progress indicators
    phases = [
        ("✅", "Fundamentos", "Completado"),
        ("✅", "Agentes Core", "Completado"),
        ("🔄", "Automatización", "En progreso"),
        ("⏳", "Dashboard", "Próximo"),
        ("⏳", "Mejoras", "Futuro"),
    ]
    
    start_x = 0.8
    for i, (icon, label, status) in enumerate(phases):
        x = start_x + i * 2.5
        
        # Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(1.3), Inches(0.6), Inches(0.6))
        if status == "Completado":
            circle.fill.solid()
            circle.fill.fore_color.rgb = GREEN_500
        elif status == "En progreso":
            circle.fill.solid()
            circle.fill.fore_color.rgb = GOLD_500
        else:
            circle.fill.solid()
            circle.fill.fore_color.rgb = NAVY_700
        circle.line.fill.background()
        
        # Icon in circle
        tf = circle.text_frame
        tf.paragraphs[0].text = icon
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Label
        add_title_text(slide, label, x - 0.3, 2.0, 1.2, 0.4, font_size=11, bold=True, color=WHITE)
        add_title_text(slide, status, x - 0.3, 2.3, 1.2, 0.3, font_size=9, bold=False, color=NAVY_100)
        
        # Connecting line
        if i < len(phases) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.6), Inches(1.55), Inches(1.9), Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = NAVY_700
            line.line.fill.background()
    
    # Cards
    add_card(slide, 0.5, 3.0, 4.0, 3.8,
        "✅ COMPLETADO",
        [
            "6 Agentes IA especializados",
            "Chat con lenguaje natural",
            "Base de datos completa",
            "Verificación KYC (Stripe)",
            "Generación contratos PDF",
            "Sistema de referidos",
            "Pagos automáticos (Stripe)",
            "Panel de clientes y propiedades"
        ],
        GREEN_500
    )
    
    add_card(slide, 4.8, 3.0, 4.0, 1.8,
        "🔄 EN PROGRESO",
        [
            "Notificaciones automáticas",
            "Recordatorios de pago",
            "Alertas de vencimiento"
        ],
        GOLD_500
    )
    
    add_card(slide, 4.8, 5.0, 4.0, 1.8,
        "📋 PRÓXIMO",
        [
            "Dashboard visual",
            "Reportes avanzados",
            "App móvil clientes"
        ],
        NAVY_100
    )
    
    # =========================================================================
    # SLIDE 3: 6 Agents
    # =========================================================================
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY_900)
    
    add_title_text(slide, "🤖 6 Agentes Inteligentes", 0.5, 0.3, 12, 0.8, font_size=36, bold=True, color=GOLD_500)
    add_title_text(slide, "Cada agente está especializado en un proceso de la cadena de valor", 0.5, 0.9, 12, 0.4, font_size=14, bold=False, color=NAVY_100)
    
    agents = [
        ("📢", "COMERCIALIZAR", "Evalúa riesgo crediticio, capta leads, gestiona recuperación"),
        ("🏠", "ADQUIRIR", "Busca propiedades, evalúa con checklist, calcula ofertas"),
        ("👤", "INCORPORAR", "Crea clientes, KYC, DTI, contratos RTO, referidos"),
        ("📊", "GESTIONAR", "Pagos automáticos, monitoreo, riesgo del portafolio"),
        ("💰", "FONDEAR", "Pipeline inversionistas, notas de deuda, compliance"),
        ("🎉", "ENTREGAR", "Elegibilidad compra, transferencia título, bonos"),
    ]
    
    for i, (icon, name, desc) in enumerate(agents):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 1.5 + row * 2.8
        
        # Card background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.9), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY_800
        shape.line.color.rgb = NAVY_700
        
        # Icon circle
        icon_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 1.4), Inches(y + 0.2), Inches(1.0), Inches(0.8))
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = GOLD_500
        icon_shape.line.fill.background()
        tf = icon_shape.text_frame
        tf.paragraphs[0].text = icon
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Name
        add_title_text(slide, name, x + 0.2, y + 1.1, 3.5, 0.4, font_size=16, bold=True, color=WHITE)
        
        # Description
        add_title_text(slide, desc, x + 0.2, y + 1.5, 3.5, 0.7, font_size=11, bold=False, color=NAVY_100)
        
        # Status badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 1.2), Inches(y + 2.1), Inches(1.5), Inches(0.3))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(220, 252, 231)  # Green-100
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].text = "✓ ACTIVO"
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = GREEN_500
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # =========================================================================
    # SLIDE 4: Stats
    # =========================================================================
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY_900)
    
    add_title_text(slide, "📈 El Sistema en Números", 0.5, 0.3, 12, 0.8, font_size=36, bold=True, color=GOLD_500)
    
    stats = [
        ("6", "Agentes IA\nEspecializados"),
        ("46", "Herramientas\nImplementadas"),
        ("15+", "Tablas en\nBase de Datos"),
        ("∞", "Conversaciones\ncon Memoria"),
    ]
    
    for i, (number, label) in enumerate(stats):
        x = 0.5 + i * 3.2
        
        # Card
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(3.0), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY_800
        shape.line.color.rgb = NAVY_700
        
        # Number
        add_title_text(slide, number, x + 0.2, 1.5, 2.6, 1.0, font_size=60, bold=True, color=GOLD_500)
        
        # Label
        add_title_text(slide, label, x + 0.2, 2.6, 2.6, 0.8, font_size=14, bold=False, color=NAVY_100)
    
    # Tech stack title
    add_title_text(slide, "🔧 Stack Tecnológico", 0.5, 3.8, 12, 0.5, font_size=20, bold=True, color=GOLD_400)
    
    tech = ["🧠 OpenAI GPT-4", "🔗 LangGraph", "🗄️ Supabase", "💳 Stripe", "⚡ FastAPI", "⚛️ Next.js", "🚀 Railway", "📄 PDF Gen"]
    
    for i, tech_item in enumerate(tech):
        col = i % 4
        row = i // 4
        x = 0.5 + col * 3.2
        y = 4.4 + row * 0.7
        
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.0), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY_800
        shape.line.color.rgb = NAVY_700
        
        tf = shape.text_frame
        tf.paragraphs[0].text = tech_item
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # =========================================================================
    # SLIDE 5: Demo
    # =========================================================================
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY_900)
    
    add_title_text(slide, "🎬 Demo en Vivo", 0.5, 0.3, 12, 0.8, font_size=36, bold=True, color=GOLD_500)
    
    # Left column - What we'll see
    add_title_text(slide, "Lo que veremos hoy:", 0.5, 1.1, 6, 0.4, font_size=18, bold=True, color=GOLD_400)
    
    add_card(slide, 0.5, 1.6, 6.0, 3.5,
        "",
        [
            "Registrar propiedad en inventario",
            "Crear perfil de cliente nuevo",
            "Calcular DTI y riesgo crediticio",
            "Generar contrato RTO (PDF)",
            "Configurar pago automático",
            "Evaluar riesgo de cartera",
            "Crear perfil de inversionista",
            "Sistema de códigos de referido"
        ],
        GOLD_500
    )
    
    # Right column - Next milestones
    add_title_text(slide, "Próximos hitos:", 6.8, 1.1, 6, 0.4, font_size=18, bold=True, color=GOLD_400)
    
    add_card(slide, 6.8, 1.6, 5.8, 1.6,
        "🔔 Automatización",
        [
            "Notificaciones por email/SMS",
            "Recordatorios de pago",
            "Flujos de cobranza automáticos"
        ],
        GOLD_500
    )
    
    add_card(slide, 6.8, 3.4, 5.8, 1.7,
        "📊 Visualización",
        [
            "Dashboard métricas en tiempo real",
            "Reportes exportables",
            "Gráficas de rendimiento"
        ],
        GOLD_500
    )
    
    # CTA
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(5.8), Inches(5.3), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD_500
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = "¿Comenzamos la demostración? 🚀"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY_900
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save
    output_path = "/Users/mariasebares/Documents/RAMA_AI/maninos-ai/docs/PRESENTACION_MANINOS_AI.pptx"
    prs.save(output_path)
    print(f"✅ Presentación guardada en: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()

